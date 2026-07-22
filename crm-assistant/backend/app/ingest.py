"""Document parsing, chunking, fact/opinion tagging, embedding and storage."""
import io
import re
from . import bedrock
from .db import db, dumps_vec


# ---------- File parsers ----------

def parse_docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def parse_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            vals = [str(c) for c in row if c is not None and str(c).strip()]
            if vals:
                parts.append(" | ".join(vals))
    return "\n".join(parts)


def parse_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        parts.append(t)
    return "\n".join(parts)


def parse_pdf(data: bytes) -> str:
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(data))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def parse_file(filename: str, data: bytes) -> str:
    name = filename.lower()
    if name.endswith(".docx"):
        return parse_docx(data)
    if name.endswith(".xlsx") or name.endswith(".xlsm"):
        return parse_xlsx(data)
    if name.endswith(".pptx"):
        return parse_pptx(data)
    if name.endswith(".pdf"):
        return parse_pdf(data)
    # txt, md, minutes, csv, and anything else -> decode as text
    try:
        return data.decode("utf-8", errors="ignore")
    except Exception:
        return ""


# ---------- Chunking ----------

def chunk_text(text: str, target=700, overlap=120):
    """Split into semantically-ish chunks by paragraphs, capped by char length."""
    text = re.sub(r"\n{3,}", "\n\n", text.strip())
    if not text:
        return []
    paras = [p.strip() for p in text.split("\n") if p.strip()]
    chunks, cur = [], ""
    for p in paras:
        if len(cur) + len(p) + 1 <= target:
            cur = f"{cur}\n{p}" if cur else p
        else:
            if cur:
                chunks.append(cur)
            if len(p) > target:
                # hard split long paragraph
                for i in range(0, len(p), target - overlap):
                    chunks.append(p[i:i + target])
                cur = ""
            else:
                cur = p
    if cur:
        chunks.append(cur)
    return chunks


# ---------- Ingest ----------

def ingest_chunks(texts, *, entity_id=None, document_id=None, interaction_id=None,
                  sensitivity=1, department="general", source_label=None,
                  created_by=None, default_label=None, image_path=None):
    """Classify (fact/opinion), embed, and store a list of text chunks."""
    texts = [t for t in texts if t and t.strip()]
    if not texts:
        return 0

    if default_label in ("fact", "opinion"):
        tags = [{"label": default_label, "confidence": 1.0, "source_person": None} for _ in texts]
    else:
        tags = bedrock.classify_fact_opinion(texts)

    vecs = bedrock.embed(texts, input_type="search_document")

    with db() as cur:
        for idx, (text, tag, vec) in enumerate(zip(texts, tags, vecs)):
            cur.execute(
                """INSERT INTO chunks
                   (entity_id, document_id, interaction_id, text, fact_or_opinion,
                    fo_confidence, source_person, source_label, sensitivity, department,
                    embedding, image_path, created_by)
                   VALUES (%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s,%s)""",
                (entity_id, document_id, interaction_id, text, tag["label"],
                 tag["confidence"], tag.get("source_person"), source_label,
                 sensitivity, department, dumps_vec(vec),
                 image_path if idx == 0 else None, created_by),
            )
    return len(texts)
